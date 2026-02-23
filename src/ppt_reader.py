import os
from pptx import Presentation

def extract_slides(ppt_path):
    prs = Presentation(ppt_path)
    slides_data = []

    for i, slide in enumerate(prs.slides):
        text_runs = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)

        slide_text = " ".join(text_runs).strip()

        if len(slide_text) > 40:
            slides_data.append({
                "pdf": os.path.basename(ppt_path),  # keep key name same
                "page": i + 1,                   # slide number
                "text": slide_text
            })

    return slides_data
